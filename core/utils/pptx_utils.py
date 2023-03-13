"""Utilities to export large number of figures to formated pptx files"""
import pptx
from pptx import Presentation
from pptx.util import Inches, Length, Pt
from os.path import join
from tqdm import tqdm
import numpy as np

def view_layout_params(pptx_path, slides_num=1, digits=3):
    pprs = Presentation(pptx_path)
    layout_dict = {}
    for shape in pprs.slides[slides_num].shapes:
        print("Object:", shape.name)
        print(type(shape))
        pos_dict = {"height" : np.around(Length(shape.height).inches, digits),
                    "width" : np.around(Length(shape.width).inches, digits),
                    "top" : np.around(Length(shape.top).inches, digits),
                    "left" : np.around(Length(shape.left).inches, digits),}
        print(pos_dict)
        layout_dict[shape.name] = pos_dict
        if hasattr(shape,"text"):
            print("Text ", shape.text)

        if isinstance(shape, pptx.shapes.picture.Picture):
            crop_dict={"crop_right" : np.around(shape.crop_right, digits),
                        "crop_left" : np.around(shape.crop_left, digits),
                        "crop_top" : np.around(shape.crop_top, digits),
                        "crop_bottom" : np.around(shape.crop_bottom, digits), }
            print(crop_dict)
            layout_dict[shape.name].update(crop_dict)
    return layout_dict


def layout_proto_evol_slide(slide, title_text, protopath, evol_figpath, manif_figpath):
    """Template script for layouting 3 figures
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.text_frame._set_font("Candara", 36, False, False)
    for k, v in {'height': 1.690, 'width': 9.313, 'top': 0.0, 'left': 3.103}.items():
        setattr(tf, k, Inches(v))
    pic0 = slide.shapes.add_picture(protopath, Inches(0.0), Inches(0.0), width=Inches(2.60))
    pic1 = slide.shapes.add_picture(evol_figpath, Inches(0.0), Inches(0.0), )
    pic2 = slide.shapes.add_picture(manif_figpath, Inches(0.0), Inches(0.0), )
    for k, v in {'height': 5.600, 'width': 6.109, 'top': 1.90, 'left': 1.624}.items():
        setattr(pic1, k, Inches(v))
    for k, v in {'height': 5.600, 'width': 5.600, 'top': 1.90, 'left': 7.733}.items():
        setattr(pic2, k, Inches(v))
